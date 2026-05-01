import asyncio
import os
import json
from aiogram import Bot, Dispatcher, types, F
import google.generativeai as genai
from pptx import Presentation

# Tokenlarni sistema o'zgaruvchilaridan olamiz (Xavfsizlik uchun)
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-pro')
bot = Bot(token=TELEGRAM_TOKEN)
dp = Dispatcher()

def create_pptx(topic, slides_data, filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "AI tomonidan yaratildi"
    for item in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = item.get("title", "Sarlavha")
        slide.placeholders[1].text = "\n".join(item.get("points", []))
    prs.save(filename)
    return filename

@dp.message(F.text == "/start")
async def start(m: types.Message):
    await m.answer("Mavzuni yuboring, slayd yasab beraman!")

@dp.message()
async def handle(m: types.Message):
    wait = await m.answer("Gemini o'ylamoqda...")
    prompt = f"Mavzu: {m.text}. 5 ta slayd uchun reja ber. FAQAT JSON: [{{'title': '...', 'points': ['...', '...']}}]"
    try:
        res = model.generate_content(prompt)
        js_str = res.text.replace('```json', '').replace('```', '').strip()
        data = json.loads(js_str)
        fname = f"{m.from_user.id}.pptx"
        create_pptx(m.text, data, fname)
        await bot.send_document(m.chat.id, types.FSInputFile(fname))
        os.remove(fname)
    except:
        await m.answer("Xato! Mavzuni aniqroq yozing.")
    await wait.delete()

async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
